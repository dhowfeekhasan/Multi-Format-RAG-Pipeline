import os
import fitz  # PyMuPDF
import docx2txt
from docx import Document

import pandas as pd
from pptx import Presentation
from PIL import Image
from io import BytesIO
import pytesseract
import pdfplumber

def process_file(file_path, document_type):
    ext = os.path.splitext(file_path)[-1].lower()
    output_folder = "extracted_data"
    os.makedirs(output_folder, exist_ok=True)
    txt_file_path = os.path.join(output_folder, os.path.basename(file_path).replace(ext, ".txt"))
    extracted_images = []

    if ext == ".pdf":
        extracted_images = extract_text_and_images_from_pdf(file_path, txt_file_path, output_folder)
    elif ext in [".pptx", ".ppt"]:
        extracted_images = extract_text_and_images_from_pptx(file_path, txt_file_path, output_folder)
    elif ext == ".xlsx":
        extracted_images = extract_text_from_excel(file_path, txt_file_path)
    elif ext == ".docx":
        extracted_images = extract_text_and_images_from_docx(file_path, txt_file_path, output_folder)
    elif ext in [".jpg", ".jpeg", ".png"]:
        extracted_images = extract_text_from_image(file_path, txt_file_path)
    else:
        print("Unsupported file format! Only PDF, DOCX, PPTX, XLSX, JPG, PNG allowed.")
        return None, None

    return txt_file_path, extracted_images

def extract_text_and_images_from_pdf(pdf_path, txt_file_path, output_folder):
    extracted_images = []
    all_text = []

    try:
        # Use pdfplumber to extract text
        with pdfplumber.open(pdf_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text:
                    all_text.append(f"\n--- Page {page_num + 1} Text ---\n{text}")

                # Extract tables
                tables = page.extract_tables()
                for table_num, table in enumerate(tables):
                    all_text.append(f"\n--- Page {page_num + 1}, Table {table_num + 1} ---")
                    for row in table:
                        all_text.append(" | ".join([str(cell) if cell else "" for cell in row]))

        # Save extracted text
        with open(txt_file_path, "w", encoding="utf-8") as f:
            f.write("\n".join(all_text))

    except Exception as e:
        print(f"Error reading PDF: {e}")
        return []

    # Extract images from PDF (using PyMuPDF)
    try:
        doc = fitz.open(pdf_path)
        for page_num, page in enumerate(doc):
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                img_pil = Image.open(BytesIO(image_bytes))
                img_filename = f"{os.path.basename(pdf_path)}_page{page_num+1}_img{img_index+1}.png"
                img_path = os.path.join(output_folder, img_filename)
                img_pil.save(img_path, format="PNG")
                extracted_images.append(img_path)
        print(f"Extracted {len(extracted_images)} images from PDF")
    except Exception as e:
        print(f"Image extraction from PDF failed: {e}")

    return extracted_images

def extract_text_and_images_from_docx(docx_path, txt_file_path, output_folder):
    doc = Document(docx_path)
    full_text = []
    extracted_images = []

    # Extract paragraphs
    for para in doc.paragraphs:
        full_text.append(para.text)

    # Extract tables
    for table_num, table in enumerate(doc.tables):
        full_text.append(f"\n### Table {table_num + 1} ###")
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            full_text.append(" | ".join(row_data))
    # Save text content
    with open(txt_file_path, "w", encoding="utf-8") as f:
        f.write("\n".join(full_text))

    # Extract embedded images (already handled via docx2txt before)
    for rel in doc.part._rels:
        rel = doc.part._rels[rel]
        if "image" in rel.target_ref:
            image_data = rel.target_part.blob
            img_filename = f"{os.path.basename(docx_path)}_img_{len(extracted_images)+1}.png"
            img_path = os.path.join(output_folder, img_filename)
            with open(img_path, "wb") as img_file:
                img_file.write(image_data)
            extracted_images.append(img_path)

    return extracted_images

import os
from pptx import Presentation
from PIL import Image
from io import BytesIO
import subprocess

def convert_ppt_to_pptx(ppt_path):
    """Convert .ppt file to .pptx using LibreOffice in headless mode"""
    pptx_path = ppt_path.replace(".ppt", ".pptx")
    try:
        subprocess.run(["libreoffice", "--headless", "--convert-to", "pptx", ppt_path], check=True)
        print(f"Converted {ppt_path} to {pptx_path}")
        return pptx_path
    except subprocess.CalledProcessError as e:
        print(f"Error converting PPT to PPTX: {e}")
        return None

def extract_text_and_images_from_pptx(ppt_path, txt_file_path, output_folder):
    """Extract text, tables, and images from both PPT and PPTX files"""
    # Check the file extension and convert if necessary
    if ppt_path.endswith(".ppt"):
        pptx_path = convert_ppt_to_pptx(ppt_path)
        if not pptx_path:
            print("Failed to convert .ppt to .pptx")
            return []
        ppt_path = pptx_path  # Use the converted .pptx file
    
    # Process the PPTX file
    presentation = Presentation(ppt_path)
    extracted_images = []
    full_text = []

    for slide_num, slide in enumerate(presentation.slides):
        full_text.append(f"\n### Slide {slide_num + 1} ###")

        for shape_num, shape in enumerate(slide.shapes):
            # Extract text from text boxes or titles
            if hasattr(shape, "text") and shape.text.strip():
                full_text.append(shape.text.strip())

            # Extract tables
            if shape.has_table:
                full_text.append(f"\n--- Table in Slide {slide_num + 1}, Shape {shape_num + 1} ---")
                table = shape.table
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    full_text.append(" | ".join(row_data))

            # Extract images
            if hasattr(shape, "image") and shape.image is not None:
                try:
                    image_bytes = shape.image.blob
                    img_pil = Image.open(BytesIO(image_bytes))
                    img_filename = f"{os.path.basename(ppt_path)}_slide{slide_num+1}_img{shape_num+1}.png"
                    img_path = os.path.join(output_folder, img_filename)
                    img_pil.save(img_path, format="PNG")
                    extracted_images.append(img_path)
                except Exception as e:
                    print(f"Error extracting image from shape: {e}")

    # Save extracted text and table data
    with open(txt_file_path, "w", encoding="utf-8") as f:
        f.write("\n".join(full_text))

    print(f"Extracted {len(extracted_images)} images from PPT/PPTX")
    return extracted_images


def extract_text_from_excel(excel_path, txt_file_path):
    xls = pd.ExcelFile(excel_path)
    with open(txt_file_path, "w", encoding="utf-8") as txt_file:
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name)
            txt_file.write(f"\n### Sheet: {sheet_name} ###\n")
            txt_file.write(df.to_string(index=False) + "\n")
    print(f"Extracted text from Excel: {txt_file_path}")
    return []        


import fitz  # PyMuPDF
from PIL import Image
import pytesseract
import cv2
import numpy as np 

def extract_text_from_image(image_path, txt_file_path):
    try:
        # Open image using PyMuPDF (fitz)
        doc = fitz.open(image_path)
        page = doc.load_page(0)
        pix = page.get_pixmap()
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

        # Convert the image to grayscale for better OCR accuracy
        img = img.convert('L')
        
        # Convert the image to numpy array for further processing
        img_cv = cv2.cvtColor(np.array(img), cv2.COLOR_RGB2BGR)
        
        # OCR to extract text with layout preservation (structure)
        extracted_text = pytesseract.image_to_string(img_cv, config='--psm 6')  # Using page segmentation mode 6

        # If OCR text extraction is poor, process further
        if not extracted_text.strip():
            print("Warning: No text extracted, preprocessing the image.")
            gray_img = cv2.cvtColor(img_cv, cv2.COLOR_BGR2GRAY)
            _, thresh_img = cv2.threshold(gray_img, 150, 255, cv2.THRESH_BINARY)
            extracted_text = pytesseract.image_to_string(thresh_img, config='--psm 6')
        
        # Save the extracted text to a file
        with open(txt_file_path, "w", encoding="utf-8") as txt_file:
            txt_file.write(extracted_text)
        
        print(f"Extracted text from image with layout preserved and saved to: {txt_file_path}")
        return [image_path]

    except FileNotFoundError:
        print(f"Error: Image file {image_path} not found.")
        return []
    except Exception as e:
        print(f"Error extracting text from image: {e}")
        return []